#!/usr/bin/env python
"""
best_book_to_ppt.py
Create a complete, well‑designed PowerPoint deck from a PDF book.

USAGE
-----
    python best_book_to_ppt.py <input.pdf> [output.pptx] \
           [--template brand.pptx] [--images]

KEY FEATURES
------------
• Hybrid LLM loop (GPT‑4o‑mini draft ➜ GPT‑4o review) for accuracy & low cost.
• Branded design via any .pptx template; otherwise uses default PowerPoint theme.
• Section covers + bullet‑with‑picture layout; images cached on disk.
• Automatic heading detection, safe chunking below 8 k tokens.
• Progress bars, error‑handled retries, zero global variables.
"""

from __future__ import annotations
import argparse, json, os, re, tempfile, time, hashlib
from pathlib import Path
from typing import List, Tuple

import pdfplumber, requests, certifi
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type
from tqdm import tqdm
from openai import OpenAI, OpenAIError
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as ALIGN

# ─────────────────────────── CONFIGURABLE CONSTANTS ──────────────────────────
MODEL_DRAFT   = "gpt-4o-mini"
MODEL_REVIEW  = "gpt-4o"
TOK_LIMIT_IN  = 7_000            # keep prompts < 8 k tokens
BULLETS_SLIDE = 6
REVIEW_LOOPS  = 1                # 0→no review; 1→draft+single review
IMAGE_SIZE    = "1024x1024"
IMAGE_CACHE   = Path(".dalle_cache")
IMAGE_CACHE.mkdir(exist_ok=True)
RETRIES_API   = 3
# ─────────────────────────── INITIALISE CLIENT ───────────────────────────────
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY", ""))

# -----------------------------------------------------------------------------
def parse_cli() -> argparse.Namespace:
    ap = argparse.ArgumentParser(description="PDF → polished PowerPoint")
    ap.add_argument("pdf", help="Input PDF book")
    ap.add_argument("pptx", nargs="?", default="book_summary.pptx",
                    help="Output deck (default: book_summary.pptx)")
    ap.add_argument("--template", help="Optional .pptx design template")
    ap.add_argument("--images", action="store_true",
                    help="Add DALL·E section & slide illustrations")
    return ap.parse_args()

# -----------------------------------------------------------------------------
def extract_sections(pdf_path: str) -> List[Tuple[str, str]]:
    """Return list of (heading, text) tuples."""
    with pdfplumber.open(pdf_path) as pdf:
        full = "\n".join(p.extract_text() or "" for p in pdf.pages)

    # Detect headings like "Chapter 1", "CHAPTER 2", "1 Introduction" …
    pat = re.compile(r"\n(?:(?:Chapter|CHAPTER)\s+\d+|(?:\n|^)\d{1,2}\s{2,}.{5,80})")
    heads = list(pat.finditer(full))
    if not heads:                         # Fallback: 8 000‑word chunks
        words = full.split()
        step = 8_000
        return [(f"Part {i+1}", " ".join(words[i*step:(i+1)*step]))
                for i in range(len(words)//step + 1)]

    heads.append(re.Match("", len(full), len(full)))        # sentinel
    sections = []
    for i in range(len(heads)-1):
        heading = heads[i].group().strip()
        body = full[heads[i].start():heads[i+1].start()].strip()
        # remove heading from body to prevent duplication
        body = body[len(heading):].lstrip()
        sections.append((heading, body))
    return sections

# -----------------------------------------------------------------------------
def token_safe_chunks(text: str, limit: int = TOK_LIMIT_IN) -> List[str]:
    """Crude split by chars (~3.6 chars ≈ 1 token)."""
    char_limit = limit * 4
    return [text[i:i+char_limit] for i in range(0, len(text), char_limit)]

# -----------------------------------------------------------------------------
@retry(
    retry=retry_if_exception_type(OpenAIError),
    wait=wait_exponential(multiplier=1, min=4, max=30),
    stop=stop_after_attempt(RETRIES_API),
)
def chat_json(model: str, system: str, user: str, timeout: int = 120) -> dict:
    rsp = client.chat.completions.create(
        model=model,
        temperature=0.25,
        timeout=timeout,
        response_format={"type": "json_object"},
        messages=[{"role": "system", "content": system},
                  {"role": "user",   "content": user}]
    )
    return json.loads(rsp.choices[0].message.content)

def draft_slides(chunk: str, heading: str) -> List[dict]:
    sys = (
        f"Return JSON {{slides:[{{title,bullets,speaker_notes}}]}}. "
        f"Each slide ≤{BULLETS_SLIDE} bullets, ≤20 words per bullet. "
        "Content must cover every clinically relevant idea in the text. "
        "Do NOT introduce extra topics."
    )
    user = f"Heading: {heading}\n\n{chunk}"
    return chat_json(MODEL_DRAFT, sys, user)["slides"]

def review_slides(chunk: str, slides: List[dict]) -> Tuple[bool, List[dict]]:
    sys = (
        "You are SlideQA‑GPT. Validate the draft slides against the source. "
        f"If every key idea is covered once and no slide exceeds {BULLETS_SLIDE} bullets "
        "return {ok:true}. Otherwise return {ok:false,slides:[corrected list]}. "
        "Keep wording concise, ≤20 words per bullet."
    )
    payload = json.dumps({"slides": slides}, ensure_ascii=False)
    user = f"Source:\n{chunk}\n\nDraft slides JSON:\n{payload}"
    rsp = chat_json(MODEL_REVIEW, sys, user)
    return rsp.get("ok", False), rsp.get("slides", slides)

def summarise_chunk(chunk: str, heading: str) -> List[dict]:
    slides = draft_slides(chunk, heading)
    for _ in range(REVIEW_LOOPS):
        ok, slides = review_slides(chunk, slides)
        if ok:
            break
    return slides

# -----------------------------------------------------------------------------
def dalle_cached(prompt: str) -> str:
    """Return path to cached PNG for a given prompt."""
    h = hashlib.sha1(prompt.encode()).hexdigest()[:16]
    path = IMAGE_CACHE / f"{h}.png"
    if path.exists():
        return str(path)
    url = client.images.generate(model="dall-e-3", prompt=prompt,
                                 n=1, size=IMAGE_SIZE).data[0].url
    data = requests.get(url, timeout=60, verify=certifi.where()).content
    path.write_bytes(data)
    return str(path)

# -----------------------------------------------------------------------------
def add_section_cover(prs: Presentation, heading: str, img_path: str | None):
    blank = prs.slide_layouts[6]          # blank
    slide = prs.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                         prs.slide_width-Inches(1),
                                         Inches(1.5))
    tf = title_box.text_frame
    tf.text = heading
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].alignment = ALIGN.LEFT
    if img_path:
        slide.shapes.add_picture(img_path, Inches(0),
                                 Inches(1.5), width=prs.slide_width)

def add_content_slide(prs: Presentation, slide_dict: dict, img_path: str | None):
    layout = prs.slide_layouts[1]         # title + content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = slide_dict["title"]

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for b in slide_dict["bullets"]:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(4)

    if slide_dict.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = slide_dict["speaker_notes"]

    # optional art on right half
    if img_path:
        pic_x = prs.slide_width * 0.55
        pic_y = Inches(1)
        pic_w = prs.slide_width - pic_x - Inches(0.2)
        slide.shapes.add_picture(img_path, pic_x, pic_y, width=pic_w)

# -----------------------------------------------------------------------------
def build_presentation(sections: List[Tuple[str, List[dict]]],
                       outfile: str,
                       template: str | None = None,
                       with_images: bool = False):
    prs = Presentation(template) if template else Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = Path(outfile).stem.replace("_", " ")
    subtitle = title_slide.placeholders[1]
    subtitle.text = "Auto‑generated concise presentation"

    for heading, slides in sections:
        cover_img = dalle_cached(f"Minimalist flat illustration. {heading}") if with_images else None
        add_section_cover(prs, heading, cover_img)

        for sd in slides:
            bullets = sd["bullets"]
            while bullets:                # split overflow
                chunk_dict = sd | {"bullets": bullets[:BULLETS_SLIDE]}
                img = (dalle_cached(f"Clean line art, {sd['title']}")
                       if with_images else None)
                add_content_slide(prs, chunk_dict, img)
                bullets = bullets[BULLETS_SLIDE:]

    prs.save(outfile)

# -----------------------------------------------------------------------------
def main():
    args = parse_cli()

    sections_raw = extract_sections(args.pdf)
    final_sections: List[Tuple[str, List[dict]]] = []

    for heading, text in tqdm(sections_raw, desc="Processing sections"):
        slides: List[dict] = []
        for chunk in token_safe_chunks(text):
            slides.extend(summarise_chunk(chunk, heading))
        final_sections.append((heading, slides))

    build_presentation(final_sections,
                       outfile=args.pptx,
                       template=args.template,
                       with_images=args.images)
    print("✔  Saved:", args.pptx)

# -----------------------------------------------------------------------------
main()
